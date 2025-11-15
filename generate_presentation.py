import collections.abc
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import numpy as np
import os

# --- Constants from the prompt ---
# Colors
ALFA_RED = RGBColor(0xE7, 0x4C, 0x3C)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_BLUE = RGBColor(0x34, 0x98, 0xDB)
GREEN = RGBColor(0x27, 0xAE, 0x60)
LIGHT_GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_COLOR = RGBColor(0x00, 0x00, 0x00)

# Hex versions for Matplotlib
ALFA_RED_HEX = "#E74C3C"
DARK_BLUE_HEX = "#2C3E50"
LIGHT_BLUE_HEX = "#3498DB"
GREEN_HEX = "#27AE60"

# Fonts
TITLE_FONT = "Montserrat"
BODY_FONT = "Open Sans"
METRIC_FONT = "Roboto Mono"

# --- Helper Functions ---

def set_slide_background(slide, color):
    """Sets the background color of a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_placeholder_icon(slide, left, top, width, height, text):
    """Adds a placeholder for an icon."""
    shape = slide.shapes.add_shape(1, left, top, width, height) # 1 is for rectangle
    shape.text = text
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

def add_placeholder_image(slide, left, top, width, height, text):
    """Adds a placeholder for an image/screenshot."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.text = text
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    line = shape.line
    line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    line.dash_style = 2 # Dashed line

def generate_bar_chart(categories, values, title, output_path):
    """Generates a styled bar chart and saves it as an image."""
    plt.style.use('seaborn-v0_8-whitegrid')
    fig, ax = plt.subplots(figsize=(8, 4))
    bars = ax.barh(categories, values, color=LIGHT_BLUE_HEX)
    ax.set_title(title, fontname=TITLE_FONT, fontsize=14, weight='bold')
    ax.set_xlabel("Market Size ($M)", fontname=BODY_FONT, fontsize=12)
    ax.tick_params(axis='y', labelsize=11)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.xaxis.grid(True, linestyle='--', which='major', color='grey', alpha=.25)
    ax.yaxis.grid(False)
    fig.tight_layout()
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()

def generate_dual_axis_chart(years, revenue, ebitda, output_path):
    """Generates a dual-axis chart for financials."""
    plt.style.use('seaborn-v0_8-whitegrid')
    fig, ax1 = plt.subplots(figsize=(8, 4))

    # Bar chart for revenue
    color = LIGHT_BLUE_HEX
    ax1.set_xlabel('Year', fontname=BODY_FONT, fontsize=12)
    ax1.set_ylabel('Revenue (mln RUB)', color=color, fontname=BODY_FONT, fontsize=12)
    ax1.bar(years, revenue, color=color, label='Revenue')
    ax1.tick_params(axis='y', labelcolor=color)

    # Line chart for EBITDA
    ax2 = ax1.twinx()
    color = ALFA_RED_HEX
    ax2.set_ylabel('EBITDA (mln RUB)', color=color, fontname=BODY_FONT, fontsize=12)
    ax2.plot(years, ebitda, color=color, marker='o', linestyle='-', label='EBITDA')
    ax2.tick_params(axis='y', labelcolor=color)

    fig.tight_layout()
    plt.title('Financial Forecast: Revenue & EBITDA', fontname=TITLE_FONT, fontsize=14, weight='bold')
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()


# --- Slide Creation Functions ---

def create_title_slide(prs):
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE_BG)

    # Placeholder for gradient background
    add_placeholder_image(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, "[Abstract gradient background]")

    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title.text = "Альфа-Аналитика: Монетизация данных и DS-экспертизы для B2B"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle.text = "Стратегическая дорожная карта создания продукта с доходом 1+ млрд руб./год"
    p = subtitle.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    add_placeholder_image(slide, Inches(8.5), Inches(6), Inches(1), Inches(0.5), "[Alfa-Bank Logo]")
    
    author = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    author.text = "Докладчик: [Имя], [Должность] | [Дата]"
    p = author.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER


def create_problem_statement_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_GRAY_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title.text = "ЗАДАЧА: Создание новой линии дохода на основе данных"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Icons
    add_placeholder_icon(slide, Inches(1), Inches(1.5), Inches(1.5), Inches(1.5), "[Icon: Data in safe]")
    add_placeholder_icon(slide, Inches(4.25), Inches(1.5), Inches(1.5), Inches(1.5), "[Icon: Money with wings]")
    add_placeholder_icon(slide, Inches(7.5), Inches(1.5), Inches(1.5), Inches(1.5), "[Icon: Competitors ahead]")

    # Metrics
    metrics = {
        "💰 Доход: 1+ млрд руб/год к 3-му году": Inches(0.5),
        "⏱️ Срок окупаемости: < 1 года": Inches(3),
        "🚀 Запуск: до 2027 года": Inches(5.5),
        "💸 CapEx: < 400 млн руб": Inches(7.5)
    }
    for text, left in metrics.items():
        tb = slide.shapes.add_textbox(left, Inches(3.5), Inches(2), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.name = METRIC_FONT
        p.font.size = Pt(18)
        p.font.bold = True

    question = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    question.text = "Как создать продукт, который монетизирует существующие данные и DS-экспертизу с минимальными вложениями?"
    p = question.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER


def create_market_analysis_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_GRAY_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title.text = "РЫНОЧНЫЙ ПОТЕНЦИАЛ И КОНКУРЕНТНЫЙ ЛАНДШАФТ"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True

    # Generate chart
    chart_path = "market_analysis_chart.png"
    categories = ['Недвижимость', 'Логистика', 'Ритейл/FMCG', 'B2B-аналитика в РФ']
    values = [150, 180, 320, 850]
    generate_bar_chart(categories, values, "Размер рынков (2025)", chart_path)
    
    slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8))
    os.remove(chart_path)

    # Competitors
    add_placeholder_image(slide, Inches(1), Inches(5.5), Inches(1.5), Inches(0.5), "[Sber Logo]")
    add_placeholder_image(slide, Inches(3), Inches(5.5), Inches(1.5), Inches(0.5), "[Tinkoff Logo]")
    
    text = slide.shapes.add_textbox(Inches(5), Inches(5.5), Inches(4.5), Inches(1.5))
    text.text = "Ритейл и FMCG — оптимальный сегмент для входа: высокая готовность платить, соответствие данным, умеренная конкуренция"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(16)

def create_financials_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_GRAY_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title.text = "ФИНАНСОВАЯ МОДЕЛЬ: 1+ МЛРД РУБ/ГОД К 3-МУ ГОДУ"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True

    # Generate chart
    chart_path = "financial_model_chart.png"
    years = ['Год 1', 'Год 2', 'Год 3']
    revenue = [150, 504, 1620]
    ebitda = [3, 114, 393]
    generate_dual_axis_chart(years, revenue, ebitda, chart_path)
    
    slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8))
    os.remove(chart_path)

    # Metrics
    metrics_text = "💸 CapEx: 360 млн руб (всего за 3 года)\n📈 CAC Payback: < 12 месяцев\n🔄 Churn Rate: 10% в год"
    tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(4), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = metrics_text
    p.font.name = BODY_FONT
    p.font.size = Pt(18)

    text = slide.shapes.add_textbox(Inches(5.5), Inches(5.5), Inches(4), Inches(1.5))
    text.text = "Окупаемость инвестиций менее 1 года, точка безубыточности в начале 2-го года"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(18)


def create_target_client_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title.text = "НАШ ЦЕЛЕВОЙ КЛИЕНТ: РИТЕЙЛЕР С 5-100 ТОЧКАМИ"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True

    add_placeholder_image(slide, Inches(0.5), Inches(1.5), Inches(3), Inches(4), "[Infographic of a person]")

    # Characteristics
    characteristics = {
        "📊 Выручка: 100 млн - 10 млрд руб/год": (Inches(4), Inches(1.5)),
        "🌍 Точки продаж: 5-100 в регионах России": (Inches(4), Inches(2.5)),
        "👤 Покупатель: Директор по развитию/маркетингу": (Inches(4), Inches(3.5)),
        "🎯 KPI: Рост выручки, оптимизация ассортимента": (Inches(4), Inches(4.5)),
    }
    for text, (left, top) in characteristics.items():
        tb = slide.shapes.add_textbox(left, top, Inches(5.5), Inches(0.75))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.name = BODY_FONT
        p.font.size = Pt(20)

    text = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.75))
    text.text = "Компании со средней цифровой зрелостью, ищущие конкурентное преимущество через данные"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

def create_mvp_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_GRAY_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.75))
    title.text = "MVP: 3 КЛЮЧЕВЫХ СЦЕНАРИЯ ДЛЯ РИТЕЙЛА"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 3 columns
    columns_data = {
        "Геоаналитика": (Inches(0.5), "[Map Screenshot]"),
        "Прогноз спроса": (Inches(3.75), "[Graph Screenshot]"),
        "Сегментация клиентов": (Inches(7.0), "[Pie Chart Screenshot]"),
    }
    for col_title, (left, img_text) in columns_data.items():
        tb = slide.shapes.add_textbox(left, Inches(1), Inches(2.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = col_title
        p.font.name = TITLE_FONT
        p.font.size = Pt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        add_placeholder_image(slide, left, Inches(1.7), Inches(2.5), Inches(3), img_text)

    text = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
    text.text = "API-first подход: интеграция в существующие системы клиента без необходимости разработки UI"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

def create_demo_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title.text = "ДЕМОНСТРАЦИЯ: ГЕОАНАЛИТИКА ДЛЯ ОПТИМИЗАЦИИ ЛОКАЦИЙ"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True

    add_placeholder_image(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4), "[Full Screenshot of Web Interface]")

    text = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
    text.text = "Простой интерфейс для быстрого принятия решений без специальных знаний"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

def create_architecture_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_GRAY_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title.text = "ТЕХНИЧЕСКАЯ АРХИТЕКТУРА: API-FIRST ПОДХОД"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True

    # Diagram
    diagram_text = "[Клиенты] ← API Gateway → [ML Models] ← Data Pipeline ← [Анонимизированные данные]"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = diagram_text
    p.font.name = METRIC_FONT
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Layers
    layers = {
        "Frontend: Streamlit (для демо), интеграция с CRM/ERP": (Inches(0.5), Inches(3)),
        "API Layer: FastAPI, аутентификация, рейт-лимиты": (Inches(0.5), Inches(4)),
        "ML Layer: Prophet, RandomForest, кластеризация": (Inches(5.5), Inches(3)),
        "Data Layer: Snowflake, агрегация и анонимизация": (Inches(5.5), Inches(4)),
    }
    for text, (left, top) in layers.items():
        tb = slide.shapes.add_textbox(left, top, Inches(4), Inches(0.75))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.name = BODY_FONT
        p.font.size = Pt(16)

    text = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.75))
    text.text = "Максимальное переиспользование существующей инфраструктуры и DS-решений Альфа-Банка"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

def create_legal_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title.text = "СООТВЕТСТВИЕ 152-ФЗ: АНОНИМИЗАЦИЯ И БЕЗОПАСНОСТЬ"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True

    # Diagram
    diagram_text = "1. Сбор → 2. Агрегирование + криптографическое обезличивание → 3. Предоставление"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = diagram_text
    p.font.name = METRIC_FONT
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Key points
    points = [
        "✅ Данные только агрегированные (минимум 100 человек в выборке)",
        "✅ Криптографическая анонимизация (MPC technology)",
        "✅ Запрет совместного хранения исходных и обезличенных данных",
    ]
    for i, point_text in enumerate(points):
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.8 + i * 0.8), Inches(8), Inches(0.75))
        p = tb.text_frame.paragraphs[0]
        p.text = point_text
        p.font.name = BODY_FONT
        p.font.size = Pt(20)

    text = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.75))
    text.text = "Строгий compliance с требованиями ЦБ РФ и Приказа Роскомнадзора №140"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

def create_biz_model_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_GRAY_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.75))
    title.text = "МОДЕЛЬ МОНЕТИЗАЦИИ: ПОДПИСКА + ENTERPRISE"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Tariff cards
    tariffs = {
        "Базовый": "150,000 руб/мес\n(3 модуля, 1000 запросов)",
        "Профессиональный": "400,000 руб/мес\n(все модули, безлимит)",
        "Enterprise": "от 1,000,000 руб/мес\n(кастомизация, SLA 99.95%)",
    }
    for i, (t_title, t_text) in enumerate(tariffs.items()):
        left = Inches(0.5 + i * 3.25)
        tb = slide.shapes.add_textbox(left, Inches(1.2), Inches(2.8), Inches(2))
        p_title = tb.text_frame.paragraphs[0]
        p_title.text = t_title
        p_title.font.name = TITLE_FONT
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_text = tb.text_frame.add_paragraph()
        p_text.text = t_text
        p_text.font.name = BODY_FONT
        p_text.font.size = Pt(16)

    add_placeholder_image(slide, Inches(1), Inches(3.5), Inches(8), Inches(2), "[Chart: Customer Growth Forecast]")

    text = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
    text.text = "Freemium-модель для привлечения SMB и enterprise-сегмент для основного дохода"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER

def create_swot_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.75))
    title.text = "SWOT-АНАЛИЗ: СИЛЬНЫЕ СТОРОНЫ И СТРАТЕГИИ МИНИМИЗАЦИИ РИСКОВ"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    quadrants = {
        "Strengths": ("🏗️ Технологическая база\n🧠 DS-экспертиза\n🏦 Репутация", Inches(0.5), Inches(1.2)),
        "Weaknesses": ("⚠️ Риск реидентификации\n💰 Ограниченный бюджет", Inches(5.5), Inches(1.2)),
        "Opportunities": ("📈 Рост SaaS-рынка\n🤝 Развитие обмена данными", Inches(0.5), Inches(3.5)),
        "Threats": ("⚔️ Высокая конкуренция\n🔄 Ужесточение регулирования", Inches(5.5), Inches(3.5)),
    }
    for q_title, (q_text, left, top) in quadrants.items():
        tb = slide.shapes.add_textbox(left, top, Inches(4), Inches(2))
        p_title = tb.text_frame.paragraphs[0]
        p_title.text = q_title
        p_title.font.name = TITLE_FONT
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_text = tb.text_frame.add_paragraph()
        p_text.text = q_text
        p_text.font.name = BODY_FONT
        p_text.font.size = Pt(16)

def create_roadmap_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_GRAY_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.75))
    title.text = "ДОРОЖНАЯ КАРТА: ОТ MVP ДО 1 МЛРД РУБ ДОХОДА"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Timeline
    timeline = "Q2 2025 ─── Q4 2025 ─── Q2 2026 ─── Q4 2026 ─── Q4 2027\n    │             │             │             │             │\nDiscovery → Development → Launch → Scaling → 1+ млрд руб"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = timeline
    p.font.name = METRIC_FONT
    p.font.size = Pt(14)

    # Milestones
    milestones = [
        "Q4 2025: Завершение юридической экспертизы и разработка API",
        "Q2 2026: Запуск MVP с 10 пилотными клиентами",
        "Q4 2026: 50 платящих клиентов, точка безубыточности",
        "Q4 2027: 450 клиентов, 1.62 млрд руб годового дохода",
    ]
    for i, text in enumerate(milestones):
        tb = slide.shapes.add_textbox(Inches(1), Inches(3 + i * 0.6), Inches(8), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = f"• {text}"
        p.font.name = BODY_FONT
        p.font.size = Pt(18)

def create_effect_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.75))
    title.text = "ЭФФЕКТ ДЛЯ КЛИЕНТОВ: КОНКРЕТНЫЕ МЕТРИКИ"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    effects = {
        "📉 -20%": "Снижение процента провальных точек продаж",
        "📊 +15%": "Рост среднего чека в оптимизированных локациях",
        "📦 -25%": "Сокращение издержек на хранение запасов",
    }
    for i, (metric, desc) in enumerate(effects.items()):
        left = Inches(0.5 + i * 3.25)
        tb_metric = slide.shapes.add_textbox(left, Inches(1.5), Inches(2.8), Inches(1.5))
        p_metric = tb_metric.text_frame.paragraphs[0]
        p_metric.text = metric
        p_metric.font.name = METRIC_FONT
        p_metric.font.size = Pt(40)
        p_metric.font.bold = True
        p_metric.font.color.rgb = GREEN
        p_metric.alignment = PP_ALIGN.CENTER

        tb_desc = slide.shapes.add_textbox(left, Inches(3), Inches(2.8), Inches(1))
        p_desc = tb_desc.text_frame.paragraphs[0]
        p_desc.text = desc
        p_desc.font.name = BODY_FONT
        p_desc.font.size = Pt(16)
        p_desc.alignment = PP_ALIGN.CENTER

    add_placeholder_image(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.75), "[Logos: Magnit, Lenta, Pyaterochka as pilot projects]")

def create_usp_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_GRAY_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.75))
    title.text = "ПОЧЕМУ ИМЕННО АЛЬФА-БАНК?"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    usps = {
        "🔒 Безопасность прежде всего": "MPC technology для анализа на зашифрованных данных",
        "🤖 Глубокая DS-экспертиза": "Готовые модели на основе 16 млн клиентов и 550k корпоративных счетов",
        "⚡️ Быстрое внедрение": "API-интеграция за 2 недели вместо месяцев разработки с нуля",
    }
    for i, (usp_title, usp_text) in enumerate(usps.items()):
        left = Inches(0.5 + i * 3.25)
        add_placeholder_icon(slide, left + Inches(1.1), Inches(1.2), Inches(0.6), Inches(0.6), "[Icon]")
        tb = slide.shapes.add_textbox(left, Inches(2), Inches(2.8), Inches(2.5))
        p_title = tb.text_frame.paragraphs[0]
        p_title.text = usp_title
        p_title.font.name = TITLE_FONT
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER
        p_text = tb.text_frame.add_paragraph()
        p_text.text = usp_text
        p_text.font.name = BODY_FONT
        p_text.font.size = Pt(16)
        p_text.alignment = PP_ALIGN.CENTER

    text = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
    text.text = "Не просто данные, а безопасная аналитика, которая работает в вашем бизнес-процессе"
    p = text.text_frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

def create_conclusion_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE_BG)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.75))
    title.text = "СЛЕДУЮЩИЕ ШАГИ: ЗАПУСК MVP В 2026 ГОДУ"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    steps = {
        "👥 Сбор команды": "DS-инженер, юрист по данным, продакт-менеджер",
        "🔐 Юридическая экспертиза": "Согласование процессов анонимизации с ЦБ РФ",
        "🚀 Разработка MVP": "API для геоаналитики и прогноза спроса",
    }
    for i, (step_title, step_text) in enumerate(steps.items()):
        left = Inches(0.5 + i * 3.25)
        add_placeholder_icon(slide, left + Inches(1.1), Inches(1.2), Inches(0.6), Inches(0.6), "[Icon]")
        tb = slide.shapes.add_textbox(left, Inches(2), Inches(2.8), Inches(2))
        p_title = tb.text_frame.paragraphs[0]
        p_title.text = step_title
        p_title.font.name = TITLE_FONT
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER
        p_text = tb.text_frame.add_paragraph()
        p_text.text = step_text
        p_text.font.name = BODY_FONT
        p_text.font.size = Pt(16)
        p_text.alignment = PP_ALIGN.CENTER

    cta = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    cta.text = "Готовы запустить пилотный проект с 3-5 клиентами в Q1 2026!"
    p = cta.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = ALFA_RED

    add_placeholder_image(slide, Inches(8), Inches(5), Inches(1.5), Inches(1.5), "[QR Code to Demo]")


# --- Main Function ---

def main():
    """Generates the entire presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625) # 16:9 aspect ratio

    print("Creating presentation...")

    create_title_slide(prs)
    create_problem_statement_slide(prs)
    create_market_analysis_slide(prs)
    create_target_client_slide(prs)
    create_mvp_slide(prs)
    create_demo_slide(prs)
    create_architecture_slide(prs)
    create_legal_slide(prs)
    create_biz_model_slide(prs)
    create_financials_slide(prs)
    create_swot_slide(prs)
    create_roadmap_slide(prs)
    create_effect_slide(prs)
    create_usp_slide(prs)
    create_conclusion_slide(prs)

    file_path = "Альфа-Аналитика_Презентация.pptx"
    prs.save(file_path)
    print(f"Presentation saved to {file_path}")
    print("\nNOTE: This presentation contains placeholders for icons, logos, and screenshots.")
    print("You will need to replace these manually with actual assets.")

if __name__ == "__main__":
    main()
